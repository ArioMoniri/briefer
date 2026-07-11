"""Turn a raw Telegram message into normalised, analysable text.

Handles: plain text, URLs (SSRF-safe fetch + readable-text extraction),
GitHub repos (README + metadata), PDFs, and images (kept as base64 for
the vision model). All network egress goes through the SSRF guard.
"""
from __future__ import annotations

import base64
import io
import logging
import os
import re
from dataclasses import dataclass, field
from typing import Any
from urllib.parse import urljoin, urlparse, urlunparse

import httpx
from bs4 import BeautifulSoup

from .security import safe_resolve, clamp
from .media import TWEET_RE, VIDEO_HOST_RE

log = logging.getLogger("briefer.enrich")

URL_RE = re.compile(r"https?://[^\s<>\"')]+", re.IGNORECASE)
_GITHUB_RE = re.compile(r"https?://github\.com/([^/\s]+)/([^/\s#?]+)", re.IGNORECASE)
_LUMA_RE = re.compile(r"https?://(?:lu\.ma|luma\.com)/[^\s]+", re.IGNORECASE)

_HEADERS = {"User-Agent": "BrieferBot/1.0 (+content-summariser)"}


@dataclass
class Attachment:
    kind: str            # "image" | "pdf" | "file" | "media"
    media_type: str
    data_b64: str = ""   # for images (vision)
    text: str = ""       # extracted text (pdf/file/transcript)
    filename: str = ""
    raw: bytes = b""     # for media (audio/video) awaiting transcription


@dataclass
class EnrichedContent:
    raw_text: str = ""
    urls: list[str] = field(default_factory=list)
    link_texts: dict[str, str] = field(default_factory=dict)  # url -> extracted text
    github_repos: list[dict[str, Any]] = field(default_factory=list)
    luma_urls: list[str] = field(default_factory=list)
    attachments: list[Attachment] = field(default_factory=list)
    notes: list[str] = field(default_factory=list)
    apply_links: list[str] = field(default_factory=list)  # apply/register buttons

    def as_source_block(self, per_source_limit: int = 6000) -> str:
        """Flatten everything into a single text block for the model."""
        parts: list[str] = []
        if self.apply_links:
            parts.append("=== APPLY / REGISTER LINKS (buttons found on the page) "
                         "===\n" + "\n".join(dict.fromkeys(self.apply_links)))
        if self.raw_text:
            parts.append("=== MESSAGE TEXT ===\n" + self.raw_text)
        for url, txt in self.link_texts.items():
            parts.append(f"=== LINK: {url} ===\n" + clamp(txt, per_source_limit))
        for repo in self.github_repos:
            parts.append(
                f"=== GITHUB REPO: {repo.get('full_name')} ===\n"
                f"Description: {repo.get('description', '')}\n"
                f"Stars: {repo.get('stars', '?')}  Language: {repo.get('language', '?')}\n"
                f"README (excerpt):\n" + clamp(repo.get("readme", ""), per_source_limit)
            )
        for att in self.attachments:
            if att.text:
                parts.append(
                    f"=== FILE: {att.filename} ({att.media_type}) ===\n"
                    + clamp(att.text, per_source_limit)
                )
            elif att.kind == "image":
                parts.append(f"=== IMAGE: {att.filename or 'photo'} (analysed via vision) ===")
        for note in self.notes:
            parts.append(f"[note] {note}")
        return "\n\n".join(parts) if parts else "(no extractable content)"


class Enricher:
    def __init__(self, max_bytes: int, tweet_extractor=None,
                 transcriber=None, enable_gallery_dl: bool = True,
                 gallery_max_images: int = 6, enable_browser: bool = False,
                 cookies_file: str = "", browser_profile_dir: str = "",
                 browser_storage_state: str = "", llm=None,
                 follow_nested_links: bool = False, max_nested_links: int = 3,
                 enable_link_guard: bool = True, link_guard_model: str = "",
                 safe_browsing_key: str = "") -> None:
        self.max_bytes = max_bytes
        self.tweets = tweet_extractor
        self.transcriber = transcriber
        self.enable_gallery_dl = enable_gallery_dl
        self.gallery_max_images = gallery_max_images
        self.enable_browser = enable_browser
        self.cookies_file = cookies_file
        self.browser_profile_dir = browser_profile_dir
        self.browser_storage_state = browser_storage_state
        self.llm = llm
        self.follow_nested_links = follow_nested_links
        self.max_nested_links = max_nested_links
        self.enable_link_guard = enable_link_guard
        self.link_guard_model = link_guard_model
        self.safe_browsing_key = safe_browsing_key
        # If egress goes through an HTTP proxy, the proxy owns DNS + policy, so
        # IP-pinning is unnecessary and can break the tunnel — pin only on
        # direct egress, where DNS-rebinding is a real concern.
        self._proxied = bool(
            os.environ.get("HTTPS_PROXY") or os.environ.get("https_proxy")
            or os.environ.get("ALL_PROXY") or os.environ.get("all_proxy")
        )

    def _try_gallery(self, url: str, content: EnrichedContent) -> int:
        """Fallback downloader for image posts → vision. Returns #images."""
        if not self.enable_gallery_dl:
            return 0
        from .media import gallery_images, guess_media_type
        imgs = gallery_images(url, self.gallery_max_images, self.max_bytes,
                              cookies_file=self.cookies_file)
        for data in imgs:
            content.attachments.append(
                make_image_attachment(data, guess_media_type(data), "post_image"))
        return len(imgs)

    def _fetch(self, url: str) -> httpx.Response | None:
        # Follow redirects manually so EVERY hop is validated with the SSRF
        # guard *before* a socket is opened. On direct egress we additionally
        # connect to a PINNED IP so a DNS-rebinding host can't resolve to a
        # public IP during validation and to 169.254.169.254 / 127.0.0.1
        # during the actual connect.
        current = url
        try:
            with httpx.Client(timeout=15, follow_redirects=False,
                              headers=_HEADERS, trust_env=True) as client:
                for _ in range(5):
                    ok, reason, ip = safe_resolve(current)
                    if not ok or not ip:
                        log.warning("Refusing to fetch %s: %s", current, reason)
                        return None
                    if self._proxied:
                        resp = client.get(current)
                    else:
                        parsed = urlparse(current)
                        host = parsed.hostname or ""
                        ip_host = f"[{ip}]" if ":" in ip else ip
                        netloc = ip_host + (f":{parsed.port}" if parsed.port else "")
                        ip_url = urlunparse(parsed._replace(netloc=netloc))
                        resp = client.get(
                            ip_url,
                            headers={"Host": host},
                            extensions={"sni_hostname": host},
                        )
                    if resp.is_redirect:
                        loc = resp.headers.get("location")
                        if not loc:
                            return None
                        current = urljoin(current, loc)
                        continue
                    if len(resp.content) > self.max_bytes:
                        log.warning("Body too large for %s (%d bytes)",
                                    current, len(resp.content))
                        return None
                    return resp
        except Exception as exc:  # noqa: BLE001
            log.warning("Fetch failed for %s: %s", url, exc)
            return None
        log.warning("Too many redirects for %s", url)
        return None

    def _apply_links(self, html: str, base_url: str) -> list[str]:
        """Find apply/register/CTA links (e.g. 'Hemen Başvur', 'Apply', a Luma
        or Typeform/Eventbrite/Google-Forms link) — these are usually buttons
        whose href never appears in the visible text."""
        soup = BeautifulSoup(html, "html.parser")
        out: list[str] = []
        text_kw = re.compile(
            r"apply|register|başvur|basvur|sign\s?up|join|kayıt|kayit|ticket|"
            r"rsvp|submit|get\s?started|katıl", re.IGNORECASE)
        href_kw = re.compile(
            r"(typeform\.com|lu\.ma|luma\.com|eventbrite\.|docs\.google\.com/forms|"
            r"forms\.gle|airtable\.com|form\.jotform|apply|register|application|"
            r"başvur|basvur|/rsvp|tally\.so)", re.IGNORECASE)
        for a in soup.find_all("a", href=True):
            href = a["href"].strip()
            if not href or href.startswith(("#", "mailto:", "javascript:", "tel:")):
                continue
            label = a.get_text(" ").strip()
            if text_kw.search(label) or href_kw.search(href):
                try:
                    full = urljoin(base_url, href)
                except Exception:  # noqa: BLE001
                    continue
                if full.startswith(("http://", "https://")):
                    out.append(full)
        return list(dict.fromkeys(out))[:8]

    def _extract_html_text(self, html: str) -> str:
        soup = BeautifulSoup(html, "html.parser")
        # Grab social/meta tags first — for login-walled pages (LinkedIn), the
        # og:title/og:description often carry the post excerpt even when the
        # body is empty.
        meta_bits: list[str] = []
        for prop in ("og:title", "og:description", "twitter:title",
                     "twitter:description", "description"):
            tag = (soup.find("meta", property=prop)
                   or soup.find("meta", attrs={"name": prop}))
            if tag and tag.get("content"):
                meta_bits.append(tag["content"].strip())
        for tag in soup(["script", "style", "noscript", "nav", "footer", "svg"]):
            tag.decompose()
        title = soup.title.string.strip() if soup.title and soup.title.string else ""
        text = " ".join(soup.get_text(" ").split())
        meta = "\n".join(dict.fromkeys(meta_bits))  # dedup, keep order
        return (f"{title}\n{meta}\n\n{text}").strip()

    def _github(self, owner: str, repo: str) -> dict[str, Any] | None:
        repo = repo.removesuffix(".git")
        api = f"https://api.github.com/repos/{owner}/{repo}"
        meta_resp = self._fetch(api)
        if not meta_resp or meta_resp.status_code != 200:
            return None
        meta = meta_resp.json()
        readme_resp = self._fetch(f"{api}/readme")
        readme = ""
        if readme_resp and readme_resp.status_code == 200:
            content = readme_resp.json().get("content", "")
            try:
                readme = base64.b64decode(content).decode("utf-8", "replace")
            except Exception:  # noqa: BLE001
                readme = ""
        return {
            "full_name": meta.get("full_name"),
            "description": meta.get("description") or "",
            "stars": meta.get("stargazers_count"),
            "language": meta.get("language"),
            "topics": meta.get("topics", []),
            "html_url": meta.get("html_url"),
            "readme": readme,
        }

    def _handle_tweet(self, url: str, content: EnrichedContent) -> bool:
        try:
            td = self.tweets.extract(url)
        except Exception as exc:  # noqa: BLE001
            log.warning("tweet extract failed for %s: %s", url, exc)
            return False
        if not td:
            return False
        rendered = td.render()
        # Transcribe attached tweet videos (own + quoted/retweeted/reply).
        for sub in (td, td.quoted, td.retweet_of, td.reply_to):
            if not sub:
                continue
            for vurl in sub.video_urls:
                if self.transcriber:
                    tr = self.transcriber.transcribe_url(vurl)
                    if tr.get("transcript"):
                        rendered += f"\n[video transcript]: {tr['transcript']}"
                    for frame in tr.get("keyframes", []):
                        content.attachments.append(make_image_attachment(
                            frame, "image/jpeg", "tweet_video_frame"))
            for purl in sub.photo_urls[:4]:
                resp = self._fetch(purl)
                if resp and resp.headers.get("content-type", "").startswith("image/"):
                    content.attachments.append(make_image_attachment(
                        resp.content, resp.headers["content-type"].split(";")[0],
                        "tweet_image"))
        content.link_texts[url] = rendered
        return True

    def _handle_video(self, url: str, content: EnrichedContent) -> bool:
        try:
            res = self.transcriber.transcribe_url(url)
        except Exception as exc:  # noqa: BLE001
            log.warning("video transcribe failed for %s: %s", url, exc)
            return False
        frames = res.get("keyframes", [])
        got_video = bool(res.get("transcript") or res.get("description") or frames)
        if not got_video:
            # Not actually a video (e.g. an Instagram photo post) — try the
            # gallery downloader so the images still reach the vision model.
            n = self._try_gallery(url, content)
            if n:
                content.link_texts[url] = (
                    f"Image post ({n} image(s) downloaded and analysed via vision).")
                return True
            # Nothing extractable here — return False so the caller falls back
            # to the normal HTML fetch + headless-browser render for this URL.
            return False
        block = f"Title: {res.get('title','')}\nUploader: {res.get('uploader','')}\n"
        if res.get("description"):
            block += "Description/caption:\n" + res["description"] + "\n"
        if res.get("transcript"):
            block += "Transcript:\n" + res["transcript"]
        else:
            block += "(no transcript — " + (res.get("note") or "unavailable") + ")"
        # Keyframes → analysed by the Anthropic multimodal model as images.
        for frame in frames:
            content.attachments.append(make_image_attachment(
                frame, "image/jpeg", "video_frame"))
        content.link_texts[url] = block
        return True

    def _transcribe_media_attachments(self, content: EnrichedContent) -> None:
        if not self.transcriber:
            return
        import tempfile
        for att in content.attachments:
            if att.kind != "media" or not att.raw:
                continue
            suffix = "." + (att.filename.rsplit(".", 1)[-1] if "." in att.filename
                            else "bin")
            try:
                with tempfile.NamedTemporaryFile(suffix=suffix, delete=True) as fh:
                    fh.write(att.raw)
                    fh.flush()
                    att.text = self.transcriber.transcribe_file(fh.name)
            except Exception as exc:  # noqa: BLE001
                log.warning("media transcription failed: %s", exc)
                att.text = "(could not transcribe uploaded media)"
            att.raw = b""  # free memory

    def enrich(self, text: str, attachments: list[Attachment]) -> EnrichedContent:
        content = EnrichedContent(raw_text=text or "", attachments=attachments)
        self._transcribe_media_attachments(content)
        urls = list(dict.fromkeys(URL_RE.findall(text or "")))
        content.urls = urls

        for url in urls:
            try:
                self._handle_url(url, content)
            except Exception as exc:  # noqa: BLE001
                # One problematic link must never sink the whole analysis.
                log.warning("enrich failed for %s: %s", url, exc)
                content.notes.append(f"Could not process {url}: {exc}")

        if self.follow_nested_links:
            try:
                self._follow_nested_links(content)
            except Exception as exc:  # noqa: BLE001
                log.warning("nested-link follow failed: %s", exc)
        return content

    def _follow_nested_links(self, content: EnrichedContent) -> None:
        """Find links INSIDE the fetched content (e.g. an article linked from a
        post), safety-check each, and fetch the safe ones one level deep."""
        from .link_safety import assess_link, is_probably_article

        # Everything we already have text for / already handled.
        known = set(content.urls) | set(content.link_texts.keys())
        primary_hosts = {urlparse(u).hostname or "" for u in content.urls}

        # Scan the extracted texts (not the user's own message) for new links.
        corpus = "\n".join(content.link_texts.values())
        candidates: list[str] = []
        for u in dict.fromkeys(URL_RE.findall(corpus)):
            u = u.rstrip(").,]\"'")
            if u in known or u in candidates:
                continue
            if TWEET_RE.match(u) or VIDEO_HOST_RE.match(u) or _GITHUB_RE.match(u):
                continue
            if is_probably_article(u, primary_hosts):
                candidates.append(u)

        for url in candidates[: self.max_nested_links]:
            ctx = _snippet_around(corpus, url)
            verdict = assess_link(
                url, ctx, llm=self.llm, guard_model=self.link_guard_model,
                safe_browsing_key=self.safe_browsing_key,
                enable_guard=self.enable_link_guard)
            if not verdict.safe:
                content.notes.append(
                    f"⚠️ Skipped linked page (unsafe): {url} — {verdict.reason}")
                continue
            if not verdict.fetch:  # safe but judged irrelevant
                content.notes.append(f"Linked page skipped (not relevant): {url}")
                continue
            resp = self._fetch(url)
            if not resp:
                rendered = self._browser_render(url)
                if rendered:
                    content.link_texts[url] = "[followed linked article]\n" + rendered
                continue
            ctype = resp.headers.get("content-type", "")
            if "application/pdf" in ctype:
                content.link_texts[url] = "[followed linked article]\n" + _pdf_to_text(resp.content)
            elif "text/html" in ctype or "text/plain" in ctype or not ctype:
                content.link_texts[url] = "[followed linked article]\n" + self._extract_html_text(resp.text)
            content.notes.append(
                f"✅ Followed linked article (safety ok): {url}"
                + (f" [{verdict.category}]" if verdict.category else ""))

    def _handle_url(self, url: str, content: EnrichedContent) -> None:
        if _LUMA_RE.match(url):
            content.luma_urls.append(url)
        # Tweets: extract the post + reply-parent + quoted/retweeted original,
        # pull in photos for vision, and transcribe any attached video.
        if self.tweets and TWEET_RE.match(url):
            if self._handle_tweet(url, content):
                return
        gh = _GITHUB_RE.match(url)
        if gh:
            repo = self._github(gh.group(1), gh.group(2))
            if repo:
                content.github_repos.append(repo)
                return  # repo readme is richer than the html page
        # Video links (YouTube/X/Vimeo/TikTok…): transcribe.
        if self.transcriber and VIDEO_HOST_RE.match(url):
            if self._handle_video(url, content):
                return
        resp = self._fetch(url)
        if not resp:
            # Plain fetch blocked/failed — try the headless browser if enabled.
            rendered = self._browser_render(url)
            if rendered:
                content.link_texts[url] = rendered
            else:
                content.notes.append(f"Could not fetch {url} (blocked or unreachable).")
            return
        ctype = resp.headers.get("content-type", "")
        if "application/pdf" in ctype:
            content.link_texts[url] = _pdf_to_text(resp.content)
        elif "text/html" in ctype or "text/plain" in ctype or not ctype:
            text = self._extract_html_text(resp.text)
            content.apply_links.extend(self._apply_links(resp.text, url))
            # Thin result (JS page / login wall) → try the headless browser.
            if len(text) < 400:
                rendered = self._browser_render(url)
                if len(rendered) > len(text):
                    text = rendered
            content.link_texts[url] = text
        else:
            content.notes.append(f"{url}: unsupported content-type {ctype}.")

    def _browser_render(self, url: str) -> str:
        if not self.enable_browser:
            return ""
        try:
            from .browser import fetch_rendered
            return fetch_rendered(
                url, cookies_file=self.cookies_file,
                profile_dir=self.browser_profile_dir,
                storage_state=self.browser_storage_state)
        except Exception as exc:  # noqa: BLE001
            log.warning("browser fallback error for %s: %s", url, exc)
            return ""


def _snippet_around(text: str, needle: str, span: int = 400) -> str:
    i = text.find(needle)
    if i == -1:
        return text[:span]
    start = max(0, i - span // 2)
    return text[start:start + span]


def _pdf_to_text(data: bytes) -> str:
    try:
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(data))
        pages = []
        for page in reader.pages[:40]:
            pages.append(page.extract_text() or "")
        return "\n".join(pages)
    except Exception as exc:  # noqa: BLE001
        log.warning("PDF parse failed: %s", exc)
        return "(could not extract PDF text)"


def make_image_attachment(data: bytes, media_type: str, filename: str = "") -> Attachment:
    return Attachment(
        kind="image",
        media_type=media_type,
        data_b64=base64.b64encode(data).decode(),
        filename=filename,
    )


def make_pdf_attachment(data: bytes, filename: str) -> Attachment:
    return Attachment(
        kind="pdf",
        media_type="application/pdf",
        text=_pdf_to_text(data),
        filename=filename,
    )


def make_text_attachment(data: bytes, filename: str, media_type: str) -> Attachment:
    return Attachment(
        kind="file",
        media_type=media_type,
        text=data.decode("utf-8", "replace"),
        filename=filename,
    )


def make_media_attachment(data: bytes, media_type: str, filename: str) -> Attachment:
    """Audio/video to be transcribed during enrichment."""
    return Attachment(kind="media", media_type=media_type, raw=data,
                      filename=filename or "media")


def make_office_attachment(data: bytes, filename: str, media_type: str) -> Attachment:
    return Attachment(kind="file", media_type=media_type,
                      text=office_to_text(data, filename), filename=filename)


def office_to_text(data: bytes, name: str) -> str:
    """Extract text from Word/PowerPoint/Excel files (lazy imports)."""
    n = name.lower()
    try:
        if n.endswith(".docx"):
            from docx import Document
            doc = Document(io.BytesIO(data))
            parts = [p.text for p in doc.paragraphs]
            for table in doc.tables:
                for row in table.rows:
                    parts.append(" | ".join(c.text for c in row.cells))
            return "\n".join(p for p in parts if p.strip())
        if n.endswith(".pptx"):
            from pptx import Presentation
            prs = Presentation(io.BytesIO(data))
            parts: list[str] = []
            for i, slide in enumerate(prs.slides, 1):
                parts.append(f"# Slide {i}")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = "".join(r.text for r in para.runs)
                            if t.strip():
                                parts.append(t)
            return "\n".join(parts)
        if n.endswith((".xlsx", ".xlsm")):
            from openpyxl import load_workbook
            wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
            parts = []
            for ws in wb.worksheets:
                parts.append(f"# {ws.title}")
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) for c in row if c is not None]
                    if cells:
                        parts.append(" | ".join(cells))
            return "\n".join(parts)
    except Exception as exc:  # noqa: BLE001
        log.warning("office parse failed for %s: %s", name, exc)
        return f"(could not extract text from {name})"
    return (f"(unsupported file type '{name}' — send .pdf, .docx, .pptx, "
            ".xlsx, or a text file)")
